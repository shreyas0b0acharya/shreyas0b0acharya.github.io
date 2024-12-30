from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()

# Define a slide layout (title and content layout)
slide_layout = presentation.slide_layouts[0]

# Slide 1: Title Slide
slide1 = presentation.slides.add_slide(slide_layout)

title = slide1.shapes.title
title.text = "Government School Rejuvenation"

title_format = title.text_frame.paragraphs[0]
title_format.font.bold = True
title_format.font.size = Pt(40)
title_format.font.color.rgb = RGBColor(0, 102, 204)  # Blue theme

title_placeholder = slide1.placeholders[1]
title_placeholder.text = "A National Service Scheme Initiative"

# Slide 2: Introduction Slide
slide2 = presentation.slides.add_slide(slide_layout)

title = slide2.shapes.title
title.text = "Introduction"

content = slide2.placeholders[1]
content.text = (
    "The National Service Scheme (NSS) aims to promote social responsibility and community engagement among students.\n"
    "As part of this vision, NSS volunteers from the Computer Science & Engineering Department at SJMIT organized a two-day event titled 'Government School Rejuvenation.'\n"
    "The event was conducted at two government schools:\n"
    "1. Sarkari Maadari Hiriya Prathamika Paata Shaale (Guddadarangavanahalli, 25th November 2024).\n"
    "2. Chinmuladri Rashtriya Proudha Shaale (Chitradurga, 28th November 2024)."
)

# Slide 3: Image Blocks Slide
slide3 = presentation.slides.add_slide(slide_layout)

slide3.shapes.title.text = "Event Highlights"

# Block for School 1 Image
left_block = Inches(1)
top_block = Inches(2)
width_block = Inches(4)
height_block = Inches(3)

school1_img_placeholder = slide3.shapes.add_picture(
    "sc1.jpg", left_block, top_block, width_block, height_block
)

school1_text = slide3.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(1))
school1_frame = school1_text.text_frame
school1_frame.text = "Sarkari Maadari Hiriya Prathamika Paata Shaale\n25th November 2024"

school1_frame.paragraphs[0].font.size = Pt(14)
school1_frame.paragraphs[0].font.bold = True
school1_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)  # Orange

# Block for School 2 Image
right_block = Inches(6)
top_block = Inches(2)

school2_img_placeholder = slide3.shapes.add_picture(
    "sc2.jpg", right_block, top_block, width_block, height_block
)

school2_text = slide3.shapes.add_textbox(Inches(6), Inches(5), Inches(4), Inches(1))
school2_frame = school2_text.text_frame
school2_frame.text = "Chinmuladri Rashtriya Proudha Shaale\n28th November 2024"

school2_frame.paragraphs[0].font.size = Pt(14)
school2_frame.paragraphs[0].font.bold = True
school2_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0)  # Orange

# Save the PowerPoint presentation
presentation.save("NSS_Government_School_Rejuvenation.pptx")

print("Presentation created successfully!")
